import os
import tempfile
from typing import Generator, Tuple
import PyPDF2
from docx import Document
import pandas as pd
from pptx import Presentation
import logging
import re

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class FileProcessor:
    def extract_text(self, file_path: str) -> str:
        """Extract text from various file formats"""
        ext = os.path.splitext(file_path)[1].lower()
        logger.info(f"Extracting text from {file_path} with extension {ext}")

        try:
            if ext == '.pdf':
                return self._extract_from_pdf(file_path)
            elif ext == '.docx':
                return self._extract_from_docx(file_path)
            elif ext == '.pptx':
                return self._extract_from_pptx(file_path)
            elif ext == '.xlsx':
                return self._extract_from_xlsx(file_path)
            elif ext == '.txt':
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            else:
                raise ValueError(f"Unsupported file extension: {ext}")
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            raise ValueError(f"Could not process file: {str(e)}")
        
    def extract_large_text(self, file_path: str, chunk_size: int = 1000) -> Generator[str, None, None]:
        """Yield text chunks from large files with smart splitting"""
        ext = os.path.splitext(file_path)[1].lower()
        logger.info(f"Extracting large text from {file_path}")
        
        try:
            if ext == '.pdf':
                yield from self._extract_smart_chunks_from_pdf(file_path, chunk_size)
            elif ext == '.docx':
                yield from self._extract_smart_chunks_from_docx(file_path, chunk_size)
            elif ext == '.pptx':
                yield from self._extract_smart_chunks_from_pptx(file_path, chunk_size)
            elif ext == '.xlsx':
                yield from self._extract_smart_chunks_from_xlsx(file_path, chunk_size)
            elif ext == '.txt':
                yield from self._extract_smart_chunks_from_txt(file_path, chunk_size)
            else:
                raise ValueError(f"Unsupported file extension: {ext}")
        except Exception as e:
            logger.error(f"Error extracting large text: {str(e)}")
            raise

    def _extract_smart_chunks_from_pdf(self, file_path: str, chunk_size: int) -> Generator[str, None, None]:
        """Yield PDF text in chunks with paragraph awareness"""
        current_chunk = ""
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text() or ""
                paragraphs = page_text.split('\n\n')
                
                for para in paragraphs:
                    if len(current_chunk) + len(para) > chunk_size:
                        if current_chunk:
                            yield current_chunk
                            current_chunk = para
                        else:
                            # Paragraph is too big, split by sentences
                            sentences = re.split(r'(?<=[.!?])\s+', para)
                            for sent in sentences:
                                if len(current_chunk) + len(sent) > chunk_size:
                                    if current_chunk:
                                        yield current_chunk
                                        current_chunk = sent
                                    else:
                                        # Sentence is too big, split by words
                                        words = sent.split()
                                        for word in words:
                                            if len(current_chunk) + len(word) > chunk_size:
                                                if current_chunk:
                                                    yield current_chunk
                                                    current_chunk = word
                                                else:
                                                    # Word is too big, split by characters
                                                    for i in range(0, len(word), chunk_size):
                                                        yield word[i:i+chunk_size]
                                            else:
                                                current_chunk += ' ' + word if current_chunk else word
                                else:
                                    current_chunk += ' ' + sent if current_chunk else sent
                    else:
                        current_chunk += '\n\n' + para if current_chunk else para
        
        if current_chunk:
            yield current_chunk

    def _extract_smart_chunks_from_docx(self, file_path: str, chunk_size: int) -> Generator[str, None, None]:
        """Yield DOCX text in chunks with paragraph awareness"""
        current_chunk = ""
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.text:
                if len(current_chunk) + len(para.text) > chunk_size:
                    if current_chunk:
                        yield current_chunk
                        current_chunk = para.text
                    else:
                        # Paragraph is too big, split by sentences
                        sentences = re.split(r'(?<=[.!?])\s+', para.text)
                        for sent in sentences:
                            if len(current_chunk) + len(sent) > chunk_size:
                                if current_chunk:
                                    yield current_chunk
                                    current_chunk = sent
                                else:
                                    # Sentence is too big, split by words
                                    words = sent.split()
                                    for word in words:
                                        if len(current_chunk) + len(word) > chunk_size:
                                            if current_chunk:
                                                yield current_chunk
                                                current_chunk = word
                                            else:
                                                # Word is too big, split by characters
                                                for i in range(0, len(word), chunk_size):
                                                    yield word[i:i+chunk_size]
                                        else:
                                            current_chunk += ' ' + word if current_chunk else word
                            else:
                                current_chunk += ' ' + sent if current_chunk else sent
                else:
                    current_chunk += '\n\n' + para.text if current_chunk else para.text
        if current_chunk:
            yield current_chunk


    def _extract_smart_chunks_from_pptx(self, file_path: str, chunk_size: int) -> Generator[str, None, None]:
        """Yield PPTX text in chunks with slide awareness"""
        current_chunk = ""
        prs = Presentation(file_path)
        
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            slide_content = "\n".join(slide_text)
            
            # Try to keep slide content together when possible
            if len(current_chunk) + len(slide_content) > chunk_size:
                if current_chunk:
                    yield current_chunk
                    current_chunk = slide_content
                else:
                    # Slide content is too big, split by shape
                    for shape_text in slide_text:
                        if len(current_chunk) + len(shape_text) > chunk_size:
                            if current_chunk:
                                yield current_chunk
                                current_chunk = shape_text
                            else:
                                # Shape text is too big, split by lines
                                lines = shape_text.split('\n')
                                for line in lines:
                                    if len(current_chunk) + len(line) > chunk_size:
                                        if current_chunk:
                                            yield current_chunk
                                            current_chunk = line
                                        else:
                                            # Line is too big, split by words
                                            words = line.split()
                                            for word in words:
                                                if len(current_chunk) + len(word) > chunk_size:
                                                    if current_chunk:
                                                        yield current_chunk
                                                        current_chunk = word
                                                    else:
                                                        # Word is too big, split by characters
                                                        for i in range(0, len(word), chunk_size):
                                                            yield word[i:i+chunk_size]
                                                else:
                                                    current_chunk += ' ' + word if current_chunk else word
                                    else:
                                        current_chunk += '\n' + line if current_chunk else line
                        else:
                            current_chunk += '\n' + shape_text if current_chunk else shape_text
            else:
                current_chunk += '\n\n' + slide_content if current_chunk else slide_content
        
        if current_chunk:
            yield current_chunk

    def _extract_smart_chunks_from_xlsx(self, file_path: str, chunk_size: int) -> Generator[str, None, None]:
        """Yield XLSX text in chunks with cell awareness"""
        current_chunk = ""
        df = pd.read_excel(file_path)
        
        # Convert dataframe to text representation
        for _, row in df.iterrows():
            row_text = " | ".join(str(cell) for cell in row.values if pd.notna(cell))
            
            if len(current_chunk) + len(row_text) > chunk_size:
                if current_chunk:
                    yield current_chunk
                    current_chunk = row_text
                else:
                    # Row is too big, split by cells
                    cells = [str(cell) for cell in row.values if pd.notna(cell)]
                    for cell in cells:
                        if len(current_chunk) + len(cell) > chunk_size:
                            if current_chunk:
                                yield current_chunk
                                current_chunk = cell
                            else:
                                # Cell is too big, split by words
                                words = cell.split()
                                for word in words:
                                    if len(current_chunk) + len(word) > chunk_size:
                                        if current_chunk:
                                            yield current_chunk
                                            current_chunk = word
                                        else:
                                            # Word is too big, split by characters
                                            for i in range(0, len(word), chunk_size):
                                                yield word[i:i+chunk_size]
                                    else:
                                        current_chunk += ' ' + word if current_chunk else word
                        else:
                            current_chunk += ' | ' + cell if current_chunk else cell
            else:
                current_chunk += '\n' + row_text if current_chunk else row_text
        
        if current_chunk:
            yield current_chunk

    def _extract_smart_chunks_from_txt(self, file_path: str, chunk_size: int) -> Generator[str, None, None]:
        """Yield text file content in chunks with paragraph awareness"""
        current_chunk = ""
        
        with open(file_path, 'r', encoding='utf-8') as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                    
                if len(current_chunk) + len(line) > chunk_size:
                    if current_chunk:
                        yield current_chunk
                        current_chunk = line
                    else:
                        # Line is too big, split by sentences
                        sentences = re.split(r'(?<=[.!?])\s+', line)
                        for sent in sentences:
                            if len(current_chunk) + len(sent) > chunk_size:
                                if current_chunk:
                                    yield current_chunk
                                    current_chunk = sent
                                else:
                                    # Sentence is too big, split by words
                                    words = sent.split()
                                    for word in words:
                                        if len(current_chunk) + len(word) > chunk_size:
                                            if current_chunk:
                                                yield current_chunk
                                                current_chunk = word
                                            else:
                                                # Word is too big, split by characters
                                                for i in range(0, len(word), chunk_size):
                                                    yield word[i:i+chunk_size]
                                        else:
                                            current_chunk += ' ' + word if current_chunk else word
                            else:
                                current_chunk += ' ' + sent if current_chunk else sent
                else:
                    current_chunk += '\n' + line if current_chunk else line
        
        if current_chunk:
            yield current_chunk

    # Similar smart chunking methods for other file types...

    def reconstruct_document(self, original_path: str, translated_text: str, target_lang: str) -> str:
        """Recreate document with translated text with improved handling"""
        ext = os.path.splitext(original_path)[1].lower()
        output_dir = tempfile.mkdtemp()
        output_filename = f"translated_{target_lang}{ext}"
        output_path = os.path.join(output_dir, output_filename)
        logger.info(f"Reconstructing document at {output_path}")

        try:
            if ext == '.pdf':
                from fpdf import FPDF
                pdf = FPDF()
                pdf.add_page()
                pdf.set_font("Arial", size=12)
                # Split text into chunks that fit on PDF pages
                for chunk in [translated_text[i:i+2000] for i in range(0, len(translated_text), 2000)]:
                    pdf.multi_cell(0, 10, txt=chunk)
                    pdf.add_page()
                pdf.output(output_path)
            elif ext == '.docx':
                doc = Document()
                for paragraph in translated_text.split('\n'):
                    doc.add_paragraph(paragraph)
                doc.save(output_path)
            elif ext == '.pptx':
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                for i, paragraph in enumerate(translated_text.split('\n')):
                    if i == 0:
                        slide.shapes.title.text = paragraph
                    else:
                        slide.placeholders[1].text += "\n" + paragraph
                prs.save(output_path)
            elif ext == '.xlsx':
                df = pd.DataFrame({"Translated Text": translated_text.split('\n')})
                df.to_excel(output_path, index=False)
            else:
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(translated_text)
            
            return output_path
        except Exception as e:
            logger.error(f"Document reconstruction error: {str(e)}")
            raise ValueError(f"Failed to reconstruct document: {str(e)}")